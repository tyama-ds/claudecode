"""
DuckDuckGo search client implementation.

Supports both the new 'ddgs' package and legacy 'duckduckgo-search' package.
"""

import csv
import io
import random
import time
import warnings
from typing import List, Optional
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup

from .base import BaseSearchClient, SearchResult, PageContent


# PDF extraction support (optional)
try:
    from PyPDF2 import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

# DOCX extraction support (optional)
try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

# XLSX extraction support (optional)
try:
    import openpyxl
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False

# PPTX extraction support (optional)
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False


class DuckDuckGoSearch(BaseSearchClient):
    """DuckDuckGo search client using ddgs or duckduckgo-search library."""

    def __init__(
        self,
        max_results: int = 10,
        timeout: int = 30,
        region: str = "wt-wt",
        safe_search: str = "moderate",
        extract_images: bool = True,
        max_images: int = 5,
        proxies: dict = None,
        verify_ssl: bool = True,
        simplify_min_results: int = 3,
        simplify_max_retries: int = 3,
        waf_mitigation: bool = True,
        per_domain_delay: float = 1.0,
    ):
        """
        Initialize DuckDuckGo search client.

        Args:
            max_results: Maximum number of search results
            timeout: Request timeout in seconds
            region: Search region (default: worldwide)
            safe_search: Safe search level (off, moderate, strict)
            extract_images: Whether to extract images from pages
            max_images: Maximum number of images to extract per page
            proxies: Proxy settings dict (e.g., {"http": "http://proxy:8080", "https": "http://proxy:8080"})
            verify_ssl: Verify SSL certificates
            simplify_min_results: Trigger query simplification when results <= this
            simplify_max_retries: Max simplification levels to try (1-3)
            waf_mitigation: Enable WAF/anti-bot mitigations for page fetching:
                a reused Session with browser-like headers, automatic retry
                with backoff on 429/5xx (honouring Retry-After), a randomized
                per-domain delay, and detection of soft blocks (challenge
                pages / 403 / 406) with one re-fetch attempt
            per_domain_delay: Base seconds to wait between fetches to the same
                domain (jittered ±50%); rapid-fire requests are a common WAF
                rate-limit trigger. Set 0 to disable the delay
        """
        super().__init__(
            max_results=max_results,
            timeout=timeout,
            extract_images=extract_images,
            max_images=max_images,
        )
        self.region = region
        self.safe_search = safe_search
        self.proxies = proxies
        self.verify_ssl = verify_ssl
        self.simplify_min_results = simplify_min_results
        self.simplify_max_retries = min(simplify_max_retries, 3)  # cap at 3 levels
        self.waf_mitigation = waf_mitigation
        self.per_domain_delay = max(0.0, per_domain_delay)
        self._ddgs = None
        self._ddgs_class = None
        self._using_new_package = False
        self._session = None
        self._last_fetch = {}  # domain -> last fetch monotonic time

    # --- WAF / anti-bot mitigation --------------------------------------

    # Statuses worth an automatic backoff+retry (transient / rate-limit)
    _RETRY_STATUS = (429, 500, 502, 503, 504)
    # Substrings that betray a soft WAF/anti-bot block on an HTTP 200/403 page
    _WAF_MARKERS = (
        "access denied", "reference #", "akamai", "bot detected",
        "please enable javascript", "verifying you are human", "are you a human",
        "captcha", "_incapsula_", "incapsula", "distil", "attention required",
        "cloudflare", "just a moment",
    )

    _BROWSER_HEADERS = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept": ("text/html,application/xhtml+xml,application/xml;q=0.9,"
                   "image/avif,image/webp,image/apng,*/*;q=0.8,"
                   "application/signed-exchange;v=b3;q=0.7"),
        "Accept-Language": "ja,en-US;q=0.9,en;q=0.8",
        # Deliberately omit 'br': requests can't always decode brotli cleanly
        "Accept-Encoding": "gzip, deflate",
        "Connection": "keep-alive",
        "Upgrade-Insecure-Requests": "1",
        "Sec-Fetch-Dest": "document",
        "Sec-Fetch-Mode": "navigate",
        "Sec-Fetch-Site": "none",
        "Sec-Fetch-User": "?1",
    }

    def _build_session(self) -> "requests.Session":
        """Session with browser headers, proxy, and backoff retry adapter."""
        from requests.adapters import HTTPAdapter
        try:
            from urllib3.util.retry import Retry
        except ImportError:  # very old urllib3
            from requests.packages.urllib3.util.retry import Retry

        s = requests.Session()
        s.headers.update(self._BROWSER_HEADERS)
        if self.proxies:
            s.proxies.update(self.proxies)
        s.verify = self.verify_ssl

        # Retry generously on WAF/rate-limit *status* responses (429/5xx,
        # honouring Retry-After) but fail fast on connection errors so a dead
        # or blocked host doesn't drag the sequential crawl (connect=1).
        retry = Retry(
            total=3,
            connect=1,
            read=2,
            status=3,
            backoff_factor=0.5,  # 0s, 0.5s, 1s, 2s between attempts
            status_forcelist=list(self._RETRY_STATUS),
            allowed_methods=frozenset(["GET", "HEAD"]),
            respect_retry_after_header=True,
            raise_on_status=False,
        )
        adapter = HTTPAdapter(max_retries=retry)
        s.mount("http://", adapter)
        s.mount("https://", adapter)
        return s

    def _get_session(self) -> "requests.Session":
        if self._session is None:
            self._session = self._build_session()
        return self._session

    def _polite_wait(self, url: str) -> None:
        """Sleep so same-domain requests aren't fired back-to-back."""
        if not self.per_domain_delay:
            return
        domain = urlparse(url).netloc.lower()
        last = self._last_fetch.get(domain)
        now = time.monotonic()
        if last is not None:
            wait = random.uniform(self.per_domain_delay * 0.5,
                                  self.per_domain_delay * 1.5)
            elapsed = now - last
            if elapsed < wait:
                time.sleep(wait - elapsed)
        self._last_fetch[domain] = time.monotonic()

    def _is_waf_blocked(self, response) -> bool:
        """Detect a soft WAF/anti-bot block (challenge page or block status)."""
        if response.status_code in (401, 403, 406, 429):
            return True
        if response.status_code == 200:
            ctype = response.headers.get("Content-Type", "").lower()
            if "html" in ctype or not ctype:
                body = (response.text or "")[:4000].lower()
                return any(m in body for m in self._WAF_MARKERS)
        return False

    def _fetch(self, url: str, headers: dict):
        """WAF-aware GET: polite delay, session+retry, one re-fetch on block."""
        if not self.waf_mitigation:
            return requests.get(
                url, headers=headers, timeout=self.timeout,
                allow_redirects=True, proxies=self.proxies, verify=self.verify_ssl,
            )

        session = self._get_session()
        self._polite_wait(url)
        response = session.get(
            url, headers=headers, timeout=self.timeout, allow_redirects=True,
        )
        # One re-fetch on a soft block: some WAFs pass a client on the second
        # hit once the connection/cookies are established.
        if self._is_waf_blocked(response):
            time.sleep(random.uniform(1.5, 3.0))
            response = session.get(
                url, headers=headers, timeout=self.timeout, allow_redirects=True,
            )
        return response

    def _get_ddgs_class(self):
        """Get the DDGS class from either ddgs or duckduckgo_search package."""
        if self._ddgs_class is not None:
            return self._ddgs_class

        # Try the new 'ddgs' package first (preferred)
        try:
            from ddgs import DDGS
            self._ddgs_class = DDGS
            self._using_new_package = True
            return DDGS
        except ImportError:
            pass

        # Fall back to 'duckduckgo_search' package
        try:
            # Suppress the deprecation warning
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                from duckduckgo_search import DDGS
            self._ddgs_class = DDGS
            self._using_new_package = False
            return DDGS
        except ImportError:
            raise ImportError(
                "Neither 'ddgs' nor 'duckduckgo-search' package is installed. "
                "Install with: pip install ddgs"
            )

    def _get_ddgs(self, force_new: bool = False):
        """Get or create DuckDuckGo search instance."""
        if self._ddgs is None or force_new:
            DDGS = self._get_ddgs_class()

            # Prepare proxy URL
            proxy_url = None
            if self.proxies:
                proxy_url = self.proxies.get("https") or self.proxies.get("http")

            # Suppress deprecation warnings during instantiation
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")

                # Try to create DDGS instance with various parameter combinations
                # Different versions support different parameters
                try:
                    if proxy_url:
                        self._ddgs = DDGS(proxy=proxy_url, timeout=self.timeout)
                    else:
                        self._ddgs = DDGS(timeout=self.timeout)
                except TypeError:
                    # Older versions might not support timeout parameter
                    try:
                        if proxy_url:
                            self._ddgs = DDGS(proxy=proxy_url)
                        else:
                            self._ddgs = DDGS()
                    except Exception as e:
                        print(f"Warning: Failed to initialize DDGS with proxy: {e}")
                        self._ddgs = DDGS()

        return self._ddgs

    def search(
        self,
        query: str,
        max_results: Optional[int] = None,
        retry_count: int = 2,
        **kwargs
    ) -> List[SearchResult]:
        """
        Perform a DuckDuckGo web search with automatic query simplification.

        When results are insufficient (≤ simplify_min_results), the query is
        progressively simplified and re-searched up to simplify_max_retries
        times. Results from all attempts are merged with URL deduplication.

        Args:
            query: The search query
            max_results: Override default max results
            retry_count: Number of retries on failure (per-attempt)
            **kwargs: Additional search parameters

        Returns:
            List of search results
        """
        max_results = max_results or self.max_results

        # Initial search
        results = self._execute_search(query, max_results, retry_count, **kwargs)

        # If we have enough results, return immediately
        if len(results) >= self.simplify_min_results:
            return results

        # Query simplification retry loop
        seen_urls = {r.url for r in results}
        last_simplified = query

        for level in range(1, self.simplify_max_retries + 1):
            simplified = self.simplify_query(query, level=level)

            # Skip if simplification didn't change the query
            if simplified == last_simplified:
                continue
            last_simplified = simplified

            print(f"[DuckDuckGo] Results insufficient ({len(results)}/{self.simplify_min_results}), "
                  f"retrying with simplified query (level {level}): {simplified}")

            new_results = self._execute_search(simplified, max_results, retry_count, **kwargs)

            # Merge new results (deduplicate by URL)
            for r in new_results:
                if r.url not in seen_urls:
                    results.append(r)
                    seen_urls.add(r.url)

            if len(results) >= self.simplify_min_results:
                print(f"[DuckDuckGo] Simplification successful: {len(results)} results after level {level}")
                break

        return results

    def _execute_search(
        self,
        query: str,
        max_results: int,
        retry_count: int = 2,
        **kwargs
    ) -> List[SearchResult]:
        """
        Execute a single DuckDuckGo search with retries on failure.

        Args:
            query: The search query
            max_results: Maximum number of results
            retry_count: Number of retries on failure
            **kwargs: Additional search parameters

        Returns:
            List of search results
        """
        print(f"[DuckDuckGo] Searching: {query}")

        for attempt in range(retry_count + 1):
            try:
                ddgs = self._get_ddgs(force_new=(attempt > 0))

                # Try to perform the search
                results = ddgs.text(
                    query,
                    region=kwargs.get("region", self.region),
                    safesearch=kwargs.get("safe_search", self.safe_search),
                    max_results=max_results,
                )

                # Convert generator/list to list
                results_list = list(results) if results else []

                if not results_list and attempt < retry_count:
                    print(f"DuckDuckGo search returned empty results, retrying... (attempt {attempt + 1})")
                    time.sleep(1)
                    continue

                search_results = []
                for result in results_list:
                    url = result.get("href", result.get("link", result.get("url", "")))
                    if url:  # Only add if we have a valid URL
                        search_results.append(SearchResult(
                            title=result.get("title", ""),
                            url=url,
                            snippet=result.get("body", result.get("snippet", result.get("description", ""))),
                            metadata={
                                "source": "duckduckgo",
                                "query": query,
                            }
                        ))

                return search_results

            except Exception as e:
                error_msg = str(e)
                print(f"DuckDuckGo search error (attempt {attempt + 1}): {error_msg}")

                if attempt < retry_count:
                    # Reset the DDGS instance for next attempt
                    self._ddgs = None
                    time.sleep(1 * (attempt + 1))  # Exponential backoff
                    continue

        # Return empty results after all retries failed
        print(f"DuckDuckGo search failed after {retry_count + 1} attempts for query: {query}")
        return []

    def search_news(
        self,
        query: str,
        max_results: Optional[int] = None,
        timelimit: Optional[str] = None,
        **kwargs
    ) -> List[SearchResult]:
        """
        Search DuckDuckGo news.

        Args:
            query: The search query
            max_results: Override default max results
            timelimit: Time limit (d: day, w: week, m: month)
            **kwargs: Additional search parameters

        Returns:
            List of news search results
        """
        ddgs = self._get_ddgs()
        max_results = max_results or self.max_results

        try:
            results = ddgs.news(
                query,
                region=kwargs.get("region", self.region),
                safesearch=kwargs.get("safe_search", self.safe_search),
                timelimit=timelimit,
                max_results=max_results,
            )

            search_results = []
            for result in results:
                search_results.append(SearchResult(
                    title=result.get("title", ""),
                    url=result.get("url", result.get("link", "")),
                    snippet=result.get("body", result.get("excerpt", "")),
                    metadata={
                        "source": "duckduckgo_news",
                        "query": query,
                        "date": result.get("date", ""),
                        "publisher": result.get("source", ""),
                    }
                ))

            return search_results

        except Exception as e:
            print(f"DuckDuckGo news search error: {e}")
            return []

    def search_images(
        self,
        query: str,
        max_results: Optional[int] = None,
        **kwargs
    ) -> List[dict]:
        """
        Search for images on DuckDuckGo.

        Args:
            query: The search query
            max_results: Override default max results
            **kwargs: Additional search parameters

        Returns:
            List of image results
        """
        ddgs = self._get_ddgs()
        max_results = max_results or self.max_images

        try:
            results = ddgs.images(
                query,
                region=kwargs.get("region", self.region),
                safesearch=kwargs.get("safe_search", self.safe_search),
                max_results=max_results,
            )

            return [
                {
                    "title": r.get("title", ""),
                    "url": r.get("image", ""),
                    "thumbnail": r.get("thumbnail", ""),
                    "source": r.get("url", ""),
                }
                for r in results
            ]

        except Exception as e:
            print(f"DuckDuckGo image search error: {e}")
            return []

    def get_page_content(
        self,
        url: str,
        extract_images: Optional[bool] = None,
        **kwargs
    ) -> PageContent:
        """
        Extract content from a URL using requests and BeautifulSoup.
        Supports both HTML pages and PDF documents.

        Args:
            url: The URL to fetch
            extract_images: Override default image extraction setting
            **kwargs: Additional parameters

        Returns:
            Extracted page content
        """
        extract_images = extract_images if extract_images is not None else self.extract_images

        # Per-request headers. The Session already carries browser-like
        # defaults; a Referer helps with WAFs that check request context.
        parsed = urlparse(url)
        headers = {
            "Referer": f"{parsed.scheme}://{parsed.netloc}/",
        }

        try:
            response = self._fetch(url, headers)

            # A soft WAF/anti-bot block that survived the re-fetch: surface it
            # clearly instead of feeding a challenge page into the pipeline.
            if self.waf_mitigation and self._is_waf_blocked(response):
                return PageContent(
                    url=url,
                    title="Blocked",
                    text_content=(
                        f"[Blocked by the site's WAF/anti-bot protection "
                        f"(HTTP {response.status_code}). This page likely needs "
                        f"a real browser; try crawl_mode='ai_crawl_selenium' or "
                        f"search_method='selenium'.]"
                    ),
                    metadata={"error": "waf_blocked",
                              "status_code": response.status_code},
                )

            response.raise_for_status()

            # Detect file type from Content-Type header and URL
            content_type = response.headers.get("Content-Type", "")
            file_type = self._get_file_type(url, content_type)

            # Skip macro-enabled files for security
            if file_type == 'macro_file':
                return PageContent(
                    url=url,
                    title="Unsupported File",
                    text_content="[Macro-enabled Office files (.docm, .xlsm, .pptm, .xlsb) are not supported for security reasons]",
                    metadata={"error": "macro_file_skipped", "content_type": content_type},
                )

            # Dispatch to appropriate extraction method
            if file_type == 'pdf':
                return self._extract_pdf_content(response, url)
            elif file_type == 'docx':
                return self._extract_docx_content(response, url)
            elif file_type == 'xlsx':
                return self._extract_xlsx_content(response, url)
            elif file_type == 'pptx':
                return self._extract_pptx_content(response, url)
            elif file_type == 'csv':
                return self._extract_csv_content(response, url)
            elif file_type in ('markdown', 'text'):
                return self._extract_text_content(response, url, file_type)

            # Continue with HTML processing (default)
            response.encoding = response.apparent_encoding

            soup = BeautifulSoup(response.text, "lxml")

            # Extract title
            title = ""
            if soup.title:
                title = soup.title.string or ""

            # Remove script and style elements
            for element in soup(["script", "style", "nav", "footer", "header", "aside"]):
                element.decompose()

            # Extract main text content
            text_content = self._extract_main_content(soup)

            # Extract images if enabled
            images = []
            if extract_images:
                images = self._extract_images(soup, url)

            # Extract links
            links = self._extract_links(soup, url)

            # Extract metadata
            metadata = self._extract_metadata(soup)

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                html_content=str(soup),
                images=images[:self.max_images],
                links=links,
                metadata=metadata,
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="Error",
                text_content=f"Failed to extract content: {str(e)}",
                metadata={"error": str(e)},
            )

    def _extract_main_content(self, soup: BeautifulSoup) -> str:
        """Extract main text content from page."""
        # Try to find main content area
        main_content = None
        for selector in ["main", "article", '[role="main"]', ".content", "#content"]:
            main_content = soup.select_one(selector)
            if main_content:
                break

        if not main_content:
            main_content = soup.body or soup

        # Get text content
        text = main_content.get_text(separator="\n", strip=True)
        return self._clean_text(text)

    def _extract_images(self, soup: BeautifulSoup, base_url: str) -> List[dict]:
        """Extract images from page."""
        from urllib.parse import urljoin

        images = []
        for img in soup.find_all("img"):
            src = img.get("src", "")
            if not src:
                continue

            # Make absolute URL
            if not src.startswith(("http://", "https://")):
                src = urljoin(base_url, src)

            # Check if valid image
            if not self._is_valid_image_url(src):
                continue

            images.append({
                "src": src,
                "alt": img.get("alt", ""),
                "title": img.get("title", ""),
            })

        return images

    def _extract_links(self, soup: BeautifulSoup, base_url: str) -> List[dict]:
        """Extract links from page."""
        from urllib.parse import urljoin

        links = []
        for a in soup.find_all("a", href=True):
            href = a.get("href", "")
            if not href or href.startswith("#"):
                continue

            # Make absolute URL
            if not href.startswith(("http://", "https://")):
                href = urljoin(base_url, href)

            links.append({
                "url": href,
                "text": a.get_text(strip=True)[:100],
            })

        return links[:50]  # Limit to 50 links

    def _is_pdf_url(self, url: str) -> bool:
        """Check if URL points to a PDF file."""
        url_lower = url.lower()
        return url_lower.endswith('.pdf') or '/pdf/' in url_lower

    def _extract_pdf_content(self, response: requests.Response, url: str) -> PageContent:
        """
        Extract text content from a PDF response.

        Args:
            response: HTTP response containing PDF binary data
            url: The source URL

        Returns:
            PageContent with extracted text
        """
        if not PDF_SUPPORT:
            return PageContent(
                url=url,
                title="PDF Document",
                text_content="[PDF content extraction requires PyPDF2. Install with: pip install PyPDF2]",
                metadata={"error": "PyPDF2 not installed", "content_type": "application/pdf"},
            )

        try:
            # Read PDF from response content
            pdf_file = io.BytesIO(response.content)
            reader = PdfReader(pdf_file)

            # Extract title from metadata if available
            title = "PDF Document"
            if reader.metadata:
                title = reader.metadata.get("/Title", "PDF Document") or "PDF Document"

            # Extract text from all pages
            text_parts = []
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[Page {page_num}]\n{page_text}")
                except Exception as e:
                    text_parts.append(f"[Page {page_num}] (extraction failed: {e})")

            text_content = "\n\n".join(text_parts)

            # Clean up the extracted text
            text_content = self._clean_text(text_content)

            # Extract metadata
            metadata = {
                "content_type": "application/pdf",
                "pages": len(reader.pages),
            }
            if reader.metadata:
                for key in ["/Author", "/Subject", "/Creator", "/Producer"]:
                    if reader.metadata.get(key):
                        metadata[key.lstrip("/")] = reader.metadata.get(key)

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                metadata=metadata,
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="PDF Document",
                text_content=f"[Failed to extract PDF content: {str(e)}]",
                metadata={"error": str(e), "content_type": "application/pdf"},
            )

    def _get_file_type(self, url: str, content_type: str = "") -> str:
        """
        Determine file type from URL extension or Content-Type header.

        Args:
            url: The URL to check
            content_type: The Content-Type header value

        Returns:
            File type string: 'pdf', 'docx', 'xlsx', 'pptx', 'csv', 'markdown', or 'html'
        """
        url_lower = url.lower()
        content_type_lower = content_type.lower()

        # Check for macro-enabled files (exclude these)
        macro_extensions = ['.docm', '.xlsm', '.pptm', '.xlsb']
        for ext in macro_extensions:
            if url_lower.endswith(ext):
                return 'macro_file'  # Will be skipped

        # PDF
        if url_lower.endswith('.pdf') or 'application/pdf' in content_type_lower:
            return 'pdf'

        # DOCX
        if url_lower.endswith('.docx') or 'application/vnd.openxmlformats-officedocument.wordprocessingml' in content_type_lower:
            return 'docx'

        # XLSX
        if url_lower.endswith('.xlsx') or 'application/vnd.openxmlformats-officedocument.spreadsheetml' in content_type_lower:
            return 'xlsx'

        # PPTX
        if url_lower.endswith('.pptx') or 'application/vnd.openxmlformats-officedocument.presentationml' in content_type_lower:
            return 'pptx'

        # CSV
        if url_lower.endswith('.csv') or 'text/csv' in content_type_lower:
            return 'csv'

        # Markdown
        if url_lower.endswith(('.md', '.rmd', '.markdown')):
            return 'markdown'

        # Plain text
        if url_lower.endswith('.txt') or 'text/plain' in content_type_lower:
            return 'text'

        # Default to HTML
        return 'html'

    def _extract_docx_content(self, response: requests.Response, url: str) -> PageContent:
        """
        Extract text content from a DOCX file.

        Args:
            response: HTTP response containing DOCX binary data
            url: The source URL

        Returns:
            PageContent with extracted text
        """
        if not DOCX_SUPPORT:
            return PageContent(
                url=url,
                title="Word Document",
                text_content="[DOCX content extraction requires python-docx. Install with: pip install python-docx]",
                metadata={"error": "python-docx not installed", "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
            )

        try:
            docx_file = io.BytesIO(response.content)
            doc = DocxDocument(docx_file)

            # Extract title from core properties if available
            title = "Word Document"
            if doc.core_properties.title:
                title = doc.core_properties.title

            # Extract text from paragraphs
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))

            text_content = "\n\n".join(text_parts)
            text_content = self._clean_text(text_content)

            # Extract metadata
            metadata = {
                "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            }
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            if doc.core_properties.created:
                metadata["created"] = str(doc.core_properties.created)

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                metadata=metadata,
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="Word Document",
                text_content=f"[Failed to extract DOCX content: {str(e)}]",
                metadata={"error": str(e), "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
            )

    def _extract_xlsx_content(self, response: requests.Response, url: str) -> PageContent:
        """
        Extract text content from an XLSX file.

        Args:
            response: HTTP response containing XLSX binary data
            url: The source URL

        Returns:
            PageContent with extracted text
        """
        if not XLSX_SUPPORT:
            return PageContent(
                url=url,
                title="Excel Spreadsheet",
                text_content="[XLSX content extraction requires openpyxl. Install with: pip install openpyxl]",
                metadata={"error": "openpyxl not installed", "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
            )

        try:
            xlsx_file = io.BytesIO(response.content)
            wb = openpyxl.load_workbook(xlsx_file, read_only=True, data_only=True)

            text_parts = []
            sheet_count = 0

            for sheet_name in wb.sheetnames:
                sheet_count += 1
                ws = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")

                rows_data = []
                for row in ws.iter_rows(values_only=True):
                    # Skip completely empty rows
                    if all(cell is None or str(cell).strip() == "" for cell in row):
                        continue
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    rows_data.append(" | ".join(row_text))

                if rows_data:
                    text_parts.append("\n".join(rows_data[:100]))  # Limit to 100 rows per sheet

            wb.close()

            text_content = "\n\n".join(text_parts)
            text_content = self._clean_text(text_content)

            # Extract title from filename
            title = url.split('/')[-1] if '/' in url else "Excel Spreadsheet"

            metadata = {
                "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "sheets": sheet_count,
            }

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                metadata=metadata,
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="Excel Spreadsheet",
                text_content=f"[Failed to extract XLSX content: {str(e)}]",
                metadata={"error": str(e), "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
            )

    def _extract_pptx_content(self, response: requests.Response, url: str) -> PageContent:
        """
        Extract text content from a PPTX file.

        Args:
            response: HTTP response containing PPTX binary data
            url: The source URL

        Returns:
            PageContent with extracted text
        """
        if not PPTX_SUPPORT:
            return PageContent(
                url=url,
                title="PowerPoint Presentation",
                text_content="[PPTX content extraction requires python-pptx. Install with: pip install python-pptx]",
                metadata={"error": "python-pptx not installed", "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
            )

        try:
            pptx_file = io.BytesIO(response.content)
            prs = Presentation(pptx_file)

            text_parts = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_texts = [f"[Slide {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_texts.append(" | ".join(row_text))

                if len(slide_texts) > 1:  # More than just the slide number
                    text_parts.append("\n".join(slide_texts))

            text_content = "\n\n".join(text_parts)
            text_content = self._clean_text(text_content)

            # Extract title
            title = "PowerPoint Presentation"
            if prs.core_properties.title:
                title = prs.core_properties.title

            metadata = {
                "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "slides": slide_count,
            }
            if prs.core_properties.author:
                metadata["author"] = prs.core_properties.author

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                metadata=metadata,
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="PowerPoint Presentation",
                text_content=f"[Failed to extract PPTX content: {str(e)}]",
                metadata={"error": str(e), "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
            )

    def _extract_csv_content(self, response: requests.Response, url: str) -> PageContent:
        """
        Extract text content from a CSV file.

        Args:
            response: HTTP response containing CSV data
            url: The source URL

        Returns:
            PageContent with extracted text
        """
        try:
            # Try different encodings
            content_text = None
            for encoding in ['utf-8', 'utf-8-sig', 'cp932', 'shift_jis', 'latin1']:
                try:
                    content_text = response.content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue

            if content_text is None:
                raise ValueError("Could not decode CSV with any known encoding")

            # Parse CSV
            csv_file = io.StringIO(content_text)
            reader = csv.reader(csv_file)

            text_parts = []
            row_count = 0
            for row in reader:
                row_count += 1
                if row_count > 200:  # Limit to 200 rows
                    text_parts.append(f"... ({row_count}+ rows total)")
                    break
                if any(cell.strip() for cell in row):
                    text_parts.append(" | ".join(row))

            text_content = "\n".join(text_parts)
            text_content = self._clean_text(text_content)

            # Extract title from filename
            title = url.split('/')[-1] if '/' in url else "CSV File"

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                metadata={
                    "content_type": "text/csv",
                    "rows": row_count,
                },
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="CSV File",
                text_content=f"[Failed to extract CSV content: {str(e)}]",
                metadata={"error": str(e), "content_type": "text/csv"},
            )

    def _extract_text_content(self, response: requests.Response, url: str, file_type: str = "text") -> PageContent:
        """
        Extract content from plain text files (txt, md, rmd, etc.).

        Args:
            response: HTTP response containing text data
            url: The source URL
            file_type: Type of file ('text', 'markdown')

        Returns:
            PageContent with extracted text
        """
        try:
            # Try different encodings
            content_text = None
            for encoding in ['utf-8', 'utf-8-sig', 'cp932', 'shift_jis', 'latin1']:
                try:
                    content_text = response.content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue

            if content_text is None:
                raise ValueError("Could not decode text with any known encoding")

            text_content = self._clean_text(content_text)

            # Extract title from filename
            title = url.split('/')[-1] if '/' in url else "Text Document"

            content_type_map = {
                "text": "text/plain",
                "markdown": "text/markdown",
            }

            return PageContent(
                url=url,
                title=title,
                text_content=text_content,
                metadata={
                    "content_type": content_type_map.get(file_type, "text/plain"),
                },
            )

        except Exception as e:
            return PageContent(
                url=url,
                title="Text Document",
                text_content=f"[Failed to extract text content: {str(e)}]",
                metadata={"error": str(e), "content_type": "text/plain"},
            )

    def _extract_metadata(self, soup: BeautifulSoup) -> dict:
        """Extract metadata from page."""
        metadata = {}

        # Extract meta tags
        for meta in soup.find_all("meta"):
            name = meta.get("name", meta.get("property", ""))
            content = meta.get("content", "")
            if name and content:
                metadata[name] = content

        return metadata

    def close(self):
        """Clean up resources."""
        self._ddgs = None
